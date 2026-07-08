# -*- coding: utf-8 -*-
"""生成 REMO_new2_AdaMaO 组会汇报 PPT（纯白背景、无动画、逐模块介绍）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- 配色（克制的深色 + 浅灰，无渐变无动画）----------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x22, 0x2A, 0x35)   # 正文深灰
TITLEBAR   = RGBColor(0x2F, 0x3B, 0x52)   # 标题栏深蓝灰
ACCENT     = RGBColor(0x2E, 0x6E, 0xA8)   # 强调蓝
SOFT       = RGBColor(0x5A, 0x63, 0x70)   # 次要灰
PILL       = RGBColor(0xE8, 0xEE, 0xF5)   # 标签浅底
PILLTXT    = RGBColor(0x2E, 0x6E, 0xA8)
LINE       = RGBColor(0xD9, 0xDF, 0xE6)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def _set_font(run, size, color=INK, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def blank_bg(slide):
    # 纯白背景
    bg = add_rect(slide, 0, 0, SW, SH, WHITE)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

def add_title_bar(slide, title, tag=None):
    bar = add_rect(slide, 0, 0, SW, Inches(1.15), TITLEBAR)
    tb = slide.shapes.add_textbox(Inches(0.55), 0, SW - Inches(1.1), Inches(1.15))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set_font(r, 26, WHITE, bold=True)
    if tag:
        pill = add_rect(slide, SW - Inches(2.7), Inches(0.32), Inches(2.15), Inches(0.5), PILL)
        pt = slide.shapes.add_textbox(SW - Inches(2.7), Inches(0.32), Inches(2.15), Inches(0.5))
        ptf = pt.text_frame; ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = tag
        _set_font(rr, 12, PILLTXT, bold=True)

def add_role(slide, role, y):
    tb = slide.shapes.add_textbox(Inches(0.6), y, SW - Inches(1.2), Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "▎" + role
    _set_font(r, 15, ACCENT, bold=True)

def add_bullets(slide, items, start_y, size=14, gap=6):
    tb = slide.shapes.add_textbox(Inches(0.7), start_y, SW - Inches(1.4), SH - start_y - Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        p.space_before = Pt(2)
        bullet = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = bullet + text
        _set_font(r, size - (1 if lvl else 0), INK if lvl == 0 else SOFT)
    return tb

def new_slide():
    s = prs.slides.add_slide(BLANK)
    blank_bg(s)
    return s

# ============================ 幻灯片 1：封面 ============================
s = new_slide()
add_rect(s, 0, 0, SW, SH, WHITE)
# 左侧竖条装饰（克制，非动画）
add_rect(s, 0, 0, Inches(0.22), SH, ACCENT)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), SW - Inches(1.8), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = "REMO_new2_AdaMaO"
_set_font(r, 44, TITLEBAR, bold=True)
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = "面向昂贵高维多目标优化的自适应关系学习算法"
_set_font(r2, 22, INK)
p3 = tf.add_paragraph(); r3 = p3.add_run()
r3.text = "各模块功能详解 · 组会汇报（初稿）"
_set_font(r3, 16, SOFT)
foot = s.shapes.add_textbox(Inches(0.9), Inches(6.4), SW - Inches(1.8), Inches(0.5))
ff = foot.text_frame; fp = ff.paragraphs[0]
fr = fp.add_run(); fr.text = "PlatEMO 框架 · 关系学习 + 运行时诊断 + 动态策略切换"
_set_font(fr, 12, SOFT)

# ============================ 幻灯片 2：算法定位（简短背景） ============================
s = new_slide()
add_title_bar(s, "一、算法定位与背景")
add_role(s, "REMO 关系学习家族的最新自适应版本", Inches(1.35))
add_bullets(s, [
    "问题场景：昂贵多目标优化（真实评估次数受限，本实验 FE=300）+ 高维多目标（目标维 M=10）",
    "继承 REMO 关系学习框架：用神经网络预测「解 A 是否优于解 B」的偏好关系，减少昂贵真实评估",
    "前作 WFG10 的不足：使用固定策略（始终置信度加权），无法按进化阶段 / 种群状态调整",
    "本版核心改动：加入运行时诊断 + 动态策略切换，让算法「自己看情况」选择训练与选择策略",
    "本文档聚焦：逐一介绍算法各模块的作用、输入/输出与彼此衔接",
], Inches(2.0), size=14, gap=8)

# ============================ 幻灯片 3：整体框架（主链四步） ============================
s = new_slide()
add_title_bar(s, "二、整体框架：主链四步", tag="总览")
add_role(s, "所有模块都挂在这条主链上，理解四步即把握全局", Inches(1.35))
steps = [
    ("① Hybrid PBI 混合分类", "对种群打 好/坏 标签，输出置信度与参考解 Ref"),
    ("② 运行时诊断", "计算 coverage（散不散）与 degeneracy（偏不偏），读种群状态"),
    ("③ 动态策略切换", "依据诊断结果，切换关系对训练模式 / 候选解选择模式"),
    ("④ 训练 + 筛选 + 评估", "训练关系模型 → AdaMaOSelection 选候选 → 真实评估 → RefSelect 环境选择"),
]
y = Inches(2.1)
for i, (h, d) in enumerate(steps):
    add_rect(s, Inches(0.7), y, Inches(0.45), Inches(0.95), ACCENT)
    nm = s.shapes.add_textbox(Inches(0.7), y, Inches(0.45), Inches(0.95))
    nmf = nm.text_frame; nmf.vertical_anchor = MSO_ANCHOR.MIDDLE
    nmp = nmf.paragraphs[0]; nmp.alignment = PP_ALIGN.CENTER
    nr = nmp.add_run(); nr.text = str(i+1); _set_font(nr, 22, WHITE, bold=True)
    hb = s.shapes.add_textbox(Inches(1.35), y, Inches(4.0), Inches(0.95))
    hbf = hb.text_frame; hbf.vertical_anchor = MSO_ANCHOR.MIDDLE
    hp = hbf.paragraphs[0]; hr = hp.add_run(); hr.text = h; _set_font(hr, 16, TITLEBAR, bold=True)
    db = s.shapes.add_textbox(Inches(5.5), y, SW - Inches(6.0), Inches(0.95))
    dbf = db.text_frame; dbf.word_wrap = True; dbf.vertical_anchor = MSO_ANCHOR.MIDDLE
    dp = dbf.paragraphs[0]; dr = dp.add_run(); dr.text = d; _set_font(dr, 13, INK)
    if i < 3:
        ar = s.shapes.add_textbox(Inches(0.78), y + Inches(0.95), Inches(0.3), Inches(0.25))
        arf = ar.text_frame; arp = arf.paragraphs[0]; arp.alignment = PP_ALIGN.CENTER
        arr = arp.add_run(); arr.text = "↓"; _set_font(arr, 14, SOFT)
    y = y + Inches(1.22)
note = s.shapes.add_textbox(Inches(0.7), y + Inches(0.05), SW - Inches(1.4), Inches(0.5))
nf = note.text_frame; np = nf.paragraphs[0]
nr2 = np.add_run(); nr2.text = "↻ 循环直至评估预算耗尽，输出累积的 Archive。"
_set_font(nr2, 13, ACCENT, bold=True)

# ============================ 模块幻灯片模板 ============================
def module_slide(idx, name, chain, role, bullets):
    s = new_slide()
    add_title_bar(s, f"三、模块 {idx}：{name}", tag=chain)
    add_role(s, role, Inches(1.35))
    add_bullets(s, bullets, Inches(2.0), size=14, gap=7)
    return s

# 模块①
module_slide("①", "Hybrid PBI 混合分类", "主链第①步",
    "把连续的多目标问题，转化为「谁比谁好」的标签样本",
    [
        "功能：对当前种群打 好/坏 标签（Catalog），同时输出每个解的置信度（confidence）与一组参考解（Ref）",
        "机制：全局参考向量 PBI + 动态参考解 PBI 两套打分融合，取头部为好类、其余并入坏类",
        "输出：Catalog / confidence / Ref —— 供后续关系对训练与候选筛选直接使用",
        "重要性：是后续所有模块的「原料来源」，没有这一步，主链直接断裂",
    ])

# 模块②
module_slide("②", "运行时诊断 RuntimeDiagnostics", "主链第②步",
    "读取种群状态，告诉算法「现在散不散、偏不偏」",
    [
        "coverage（覆盖率）：解的方向覆盖了多少个均匀参考向量（越高表示种群越分散）",
        "degeneracy（退化度）：目标矩阵做 SVD 后，解释 90% 能量所需的秩（越高表示挤在低维子空间）",
        "用途：coverage 控制探索强度与评估批量；degeneracy 当前仅用于 indicator 触发（计划删除）",
        "意义：这是 AdaMaO 区别于「固定策略」前作（WFG10）的关键新增模块",
    ])

# 模块③上：关系对训练三模式
module_slide("③", "动态策略切换（上）· 关系对训练三模式", "主链第③步",
    "根据模型精度与置信度，选择不同的关系对生成方式",
    [
        "conservative（默认）：使用原始关系对 GetRelationPairs，无权重，最稳",
        "curriculum（上代误差大时）：只保留置信度前 80% 样本，过滤低置信噪声（课程学习思想）",
        "weighted（模型准 + 置信高 + 覆盖低）：置信度加权 GetRelationPairs_confidence，让靠谱样本影响更大",
        "本质：一个自适应开关，而非三个独立算法；阈值正在内生化，由问题规模自动决定",
    ])

# 模块③下：候选解选择三模式
module_slide("③", "动态策略切换（下）· 候选解选择三模式", "主链第③步",
    "根据模型精度与种群状态，选择不同的候选筛选方式",
    [
        "conservative（默认）：仅按关系得分，选 n_min 个候选去真实评估",
        "explore（模型准 + 覆盖低）：关系得分 + 不确定性 + 多样性，鼓励探索空白区域",
        "indicator（种群退化高，计划删除）：用 PIEA 风格指标 SVR 对候选重排序",
        "切换依据：模型测试误差 p_err、覆盖率 coverage、退化度 degeneracy",
    ])

# 模块④
module_slide("④", "关系预测模型训练 TrainRelationModel", "主链第④步 · 训练侧",
    "学习「好/坏映射」的神经网络，并回灌误差给策略切换",
    [
        "输入：关系对样本（两个解的决策变量拼接）+ 标签 {-1, 0, +1}，转为 one-hot 编码",
        "结构：三层前馈 patternnet（节点数 1.5×Dim / Dim / 0.5×Dim），softmax 输出",
        "训练：mapminmax 归一化；weighted 模式按置信度加权（下限截断 w_min 防训练不稳定）",
        "评估：测试集分类错误率 p_err，作为下一轮策略切换的判断依据",
    ])

# 模块⑤
module_slide("⑤", "AdaMaOSelection 候选筛选", "主链第④步 · 选择侧",
    "在候选池里挑出少量解，送去昂贵的真实评估",
    [
        "流程：内层 GA 生成候选池 → model_select 用关系模型打分 → 按候选模式筛选",
        "conservative：按关系得分取前 n_min 个",
        "explore：得分 + 不确定性 + 贪心多样性（max-min distance），避免聚簇",
        "indicator：关系得分粗筛 + SVR 指标重排序（计划删除）",
        "兜底：若筛选失败，退回普通 GA 生成备选候选，保证每轮都有评估",
    ])

# 模块⑥
module_slide("⑥", "RefSelect 环境选择 + Archive", "主链第④步 · 收尾",
    "从累积的全部真实评估解中，选出下一代种群",
    [
        "Archive：累积所有真实评估过的解，是算法的最终输出",
        "RefSelect：采用 RSEA 雷达网格策略，将高维目标映射到 2D 雷达网格做环境选择",
        "输出：规模为 N 的下一代 Population，进入下一轮主链循环",
        "作用：在有限预算下维持种群的收敛性与分布性",
    ])

# 模块补充：PIEA 指标子系统
module_slide("＋", "PIEA 指标子系统（补充）", "第③步 indicator 分支",
    "轮盘选择一种指标评估种群——已证无益，计划移除",
    [
        "组成：IndicatorSelector（轮盘选 SDE / I_epsilon+ / Minkowski）、Shape_Estimate（估计 PF 形状 Lp）",
        "配套：fitrsvm 指标模型、UpdateInformation（轮盘概率反馈）、NDSort_SDR（强支配反馈）",
        "现状：消融 5 变体无统计显著差异（Friedman p=0.84），且在断开 PF（如 DTLZ7）上会塌缩",
        "处置：计划从最终版本 AdaMaO-S 彻底移除，使论文贡献更干净、更原创",
    ])

# ============================ 结尾：总结与下一步 ============================
s = new_slide()
add_title_bar(s, "四、总结与下一步")
add_role(s, "模块虽多，但绝大多数挂在②③的增强上；已用消融定位无用模块", Inches(1.35))
add_bullets(s, [
    "主链闭环：分类 → 诊断 → 动态切策略 → 训练筛选，实现昂贵高维多目标下的自适应",
    "各模块职责清晰：分类提供标签、诊断读状态、切换选策略、训练与筛选落地、环境选择维持种群",
    "指标子系统与 curriculum 等经消融验证非必需，正在按结果裁剪",
    "下一步：删 indicator / curriculum、把超参数内生化 → 推出更干净的 AdaMaO-S（基本无外部调参）",
], Inches(2.0), size=14, gap=8)

out = r"d:\PlatEMO-master\PlatEMO-master\PlatEMO\Algorithms\Multi-objective optimization\REMO_new2_AdaMaO\REMO_new2_AdaMaO_组会汇报PPT.pptx"
prs.save(out)
print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))
